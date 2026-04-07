"""
create_pptx.py - Generates the project presentation PPTX
CSCI 3010U - Simulation and Modeling
Yash Patel - 100785833

Run with:  python create_pptx.py
Requires:  pip install python-pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────────────────────────────────────────────
# Font pairing  (visually distinct title vs body)
# ─────────────────────────────────────────────────────────────────────────────
FONT_TITLE = "Georgia"        # serif — used for all slide titles
FONT_BODY  = "Calibri Light"  # sans-serif — used for all body / bullet text
FONT_CODE  = "Courier New"    # monospace — used for formulas

# ─────────────────────────────────────────────────────────────────────────────
# Colour palette
# ─────────────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x0F, 0x1A)
PANEL       = RGBColor(0x10, 0x16, 0x2A)
BLUE_ACCENT = RGBColor(0x5B, 0x9B, 0xFF)
GREEN       = RGBColor(0x37, 0xD3, 0x5F)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT       = RGBColor(0xC8, 0xDE, 0xFF)
DIM         = RGBColor(0x6C, 0x80, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name=FONT_BODY, font_size=Pt(18),
                bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Single-run textbox."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.name       = font_name
    run.font.size       = font_size
    run.font.bold       = bold
    run.font.color.rgb  = color
    return tb


def footer(slide):
    """Standard footer bar on every slide."""
    add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35),
             fill_color=PANEL)
    add_textbox(slide,
                "CSCI 3010U  |  Yash Patel  |  100785833",
                Inches(0.4), SLIDE_H - Inches(0.34),
                Inches(8), Inches(0.32),
                font_name=FONT_BODY, font_size=Pt(11), color=DIM)


def header_strip(slide, title_text):
    """Top title bar used on content slides."""
    add_rect(slide, 0, 0,           SLIDE_W, Inches(0.08), fill_color=BLUE_ACCENT)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.12), fill_color=PANEL)
    # Title uses FONT_TITLE (serif) — clearly different from body
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.14),
                Inches(12.5), Inches(1.0),
                font_name=FONT_TITLE, font_size=Pt(34),
                bold=True, color=BLUE_ACCENT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def bullet_slide(prs, title, bullets):
    """
    Bullet slide layout.
    bullets: list of (text, level)
      level 0  = main point  — FONT_TITLE bold, 21 pt, LIGHT colour
      level 1  = sub-point   — FONT_BODY,       18 pt, DIM colour
      level -1 = formula     — FONT_CODE,        17 pt, GREEN colour
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    header_strip(slide, title)
    footer(slide)

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.45),
                                  Inches(12.2), Inches(5.75))
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for (text, level) in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        run = p.add_run()
        run.text = text

        if level == 0:                          # main heading line
            run.font.name      = FONT_TITLE     # Georgia (serif)
            run.font.size      = Pt(21)
            run.font.bold      = True
            run.font.color.rgb = LIGHT
            p.space_before     = Pt(14)
        elif level == -1:                       # formula / code
            run.font.name      = FONT_CODE      # Courier New
            run.font.size      = Pt(17)
            run.font.bold      = False
            run.font.color.rgb = GREEN
            p.space_before     = Pt(4)
            p.level            = 1
        else:                                   # sub-bullet (level 1)
            run.font.name      = FONT_BODY      # Calibri Light (sans)
            run.font.size      = Pt(18)
            run.font.bold      = False
            run.font.color.rgb = DIM
            p.space_before     = Pt(3)
            p.level            = 1

    return slide


def image_slide(prs, title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    header_strip(slide, title)
    footer(slide)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path,
                                 Inches(0.5), Inches(1.35),
                                 Inches(12.3), Inches(5.75))
    if caption:
        add_textbox(slide, caption,
                    Inches(0.5), SLIDE_H - Inches(0.62),
                    Inches(12.3), Inches(0.35),
                    font_name=FONT_BODY, font_size=Pt(12),
                    color=DIM, align=PP_ALIGN.CENTER)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_ACCENT)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12),
             fill_color=GREEN)

    # Centre card
    add_rect(slide, Inches(1.1), Inches(1.5), Inches(11.1), Inches(4.5),
             fill_color=PANEL, line_color=BLUE_ACCENT, line_width=Pt(1.8))

    # Decorative coloured ball circles
    ball_cols = [(0xDA,0x34,0x34),(0x34,0x94,0xDE),(0xCD,0xAA,0x1C),
                 (0x41,0xC0,0x41),(0xC0,0x41,0xC0),(0xDE,0x82,0x28)]
    for i, (r, g, b) in enumerate(ball_cols):
        cx = int(Inches(2.0 + i * 1.55)); cy = int(Inches(2.1)); cr = int(Inches(0.3))
        s = slide.shapes.add_shape(9, cx-cr, cy-cr, cr*2, cr*2)
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor(r,g,b)
        s.line.fill.background()

    # Main title — FONT_TITLE (Georgia, serif)
    add_textbox(slide, "2D Billiards Physics Simulation",
                Inches(1.3), Inches(2.5), Inches(10.7), Inches(1.3),
                font_name=FONT_TITLE, font_size=Pt(46),
                bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle — FONT_BODY (Calibri Light, sans)
    add_textbox(slide,
                "Simulating How Billiard Balls Move and Collide Using Python",
                Inches(1.3), Inches(3.75), Inches(10.7), Inches(0.6),
                font_name=FONT_BODY, font_size=Pt(20),
                color=LIGHT, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Yash Patel  |  100785833  |  yash.patel@ontariotechu.net",
                Inches(1.3), Inches(4.42), Inches(10.7), Inches(0.45),
                font_name=FONT_BODY, font_size=Pt(16),
                color=DIM, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "CSCI 3010U – Simulation and Modeling  |  Ontario Tech University  |  Winter 2026",
                Inches(1.3), Inches(4.88), Inches(10.7), Inches(0.45),
                font_name=FONT_BODY, font_size=Pt(14),
                color=DIM, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2 — What Is This Project?
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "What Is This Project?",
        [
            ("A computer simulation of 8 billiard balls on a 2D table", 0),
            ("Each ball moves, bounces off walls, and collides with other balls", 1),
            ("Built on real physics principles — not just animation", 1),
            ("", 0),
            ("Why billiards?", 0),
            ("A clear, testable example of a multi-body physics system", 1),
            ("Easy to observe energy and momentum conservation in action", 1),
            ("", 0),
            ("What can you do with it?", 0),
            ("Watch 8 balls bounce in real-time at 120 FPS", 1),
            ("Adjust restitution — controls how bouncy collisions are", 1),
            ("Adjust damping — controls how much friction slows the balls", 1),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3 — How Balls Move (Euler Method)
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "How Balls Move — The Euler Method",
        [
            ("Each ball has:  position (x, y)  ·  velocity (vx, vy)  ·  mass  ·  radius", 0),
            ("", 0),
            ("Euler Method — update position every tiny time step dt", 0),
            ("x(t + dt)  =  x(t)  +  vx × dt", -1),
            ("y(t + dt)  =  y(t)  +  vy × dt", -1),
            ("", 0),
            ("Damping — simulates friction on the table", 0),
            ("velocity  =  velocity × (1 − λ × dt)", -1),
            ("Higher λ  →  balls slow down faster", 1),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4 — Ball–Wall Collisions
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "Ball–Wall Collisions",
        [
            ("When a ball hits the edge of the table, it bounces back", 0),
            ("Hit left or right wall  →  vx becomes −vx", 1),
            ("Hit top or bottom wall  →  vy becomes −vy", 1),
            ("", 0),
            ("Position correction — ball is nudged back inside the boundary", 0),
            ("Prevents balls from tunnelling through the wall", 1),
            ("", 0),
            ("Wall collisions are always elastic — no energy is lost at walls", 0),
            ("Energy can only be lost during ball-ball collisions (via restitution)", 1),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5 — Ball–Ball Collisions
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "Ball–Ball Collisions — The Physics",
        [
            ("Step 1 — Detect:  distance between centres < r1 + r2", 0),
            ("", 0),
            ("Step 2 — Respond:  impulse using Newton's Law of Restitution", 0),
            ("J  =  −(1 + e) × dvn  ÷  (1/m1 + 1/m2)", -1),
            ("e = restitution (elasticity parameter)  ·  dvn = relative speed along normal", 1),
            ("", 0),
            ("Step 3 — Separate:  push balls apart to remove overlap", 0),
            ("", 0),
            ("Restitution = 1.0  →  Perfectly elastic  (no kinetic energy lost)", 0),
            ("Restitution = 0.0  →  Perfectly inelastic  (maximum energy lost)", 1),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6 — Live Demo screenshot
    # ══════════════════════════════════════════════════════════════════════════
    image_slide(prs,
        "The Interactive Simulation (Live Demo)",
        "simulation_screenshot.png",
        caption="Controls:  ↑↓ change damping  |  ←→ change restitution  |  SPACE pauses  |  R resets  |  Running at 120 FPS")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7 — Experiment 1: Restitution
    # ══════════════════════════════════════════════════════════════════════════
    image_slide(prs,
        "Experiment 1 — Effect of Restitution on Kinetic Energy",
        "experiment1_restitution.png",
        caption="Restitution = 1.0 → energy stays constant  |  Lower restitution → energy drops at each collision")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8 — Experiment 2: Damping
    # ══════════════════════════════════════════════════════════════════════════
    image_slide(prs,
        "Experiment 2 — Effect of Damping (Friction) on Kinetic Energy",
        "experiment2_damping.png",
        caption="λ = 0 → energy conserved  |  λ = 0.2 → gradual decay  |  λ = 0.6 → balls stop within ~4 seconds")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9 — Experiment 3: Momentum
    # ══════════════════════════════════════════════════════════════════════════
    image_slide(prs,
        "Experiment 3 — Total Momentum Over Time",
        "experiment3_momentum.png",
        caption="Ball-ball collisions conserve momentum in the absence of external forces. Wall reflections change the system momentum.")

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10 — Key Findings
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "Key Findings",
        [
            ("Restitution (elasticity parameter):", 0),
            ("Restitution = 1.0  →  kinetic energy stays perfectly constant", 1),
            ("Restitution = 0.7  →  energy drops in steps at each collision", 1),
            ("Restitution = 0.4  →  energy drops much faster after every hit", 1),
            ("", 0),
            ("Damping (friction coefficient):", 0),
            ("λ = 0.0  →  balls never stop, energy is fully conserved", 1),
            ("λ = 0.2  →  balls slow down gradually over time", 1),
            ("λ = 0.6  →  balls come to rest within approximately 4 seconds", 1),
            ("", 0),
            ("Momentum:", 0),
            ("Ball-ball collisions conserve momentum in the absence of external forces", 1),
            ("Wall reflections change the system momentum, as expected", 1),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 11 — Project Goals Achieved
    # ══════════════════════════════════════════════════════════════════════════
    bullet_slide(prs, "Project Goals — All Achieved  ✓",
        [
            ("✓  Euler method implemented and working correctly", 0),
            ("Position updates match expected numerical integration output", 1),
            ("", 0),
            ("✓  Realistic wall bouncing via elastic velocity reversal", 0),
            ("", 0),
            ("✓  Ball-ball collisions using Newton's Law of Restitution", 0),
            ("Energy and momentum formulas verified by all three experiments", 1),
            ("", 0),
            ("✓  Damping (friction) produces exponential energy decay", 0),
            ("Clearly confirmed in Experiment 2 graphs", 1),
            ("", 0),
            ("✓  All 3 experiments completed — graphs match expected physics", 0),
        ])

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 12 — Conclusion
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=BLUE_ACCENT)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), fill_color=GREEN)
    add_rect(slide, Inches(1.1), Inches(1.1), Inches(11.1), Inches(5.3),
             fill_color=PANEL, line_color=BLUE_ACCENT, line_width=Pt(1.8))

    # Title — FONT_TITLE (serif)
    add_textbox(slide, "Conclusion",
                Inches(1.3), Inches(1.2), Inches(10.7), Inches(0.85),
                font_name=FONT_TITLE, font_size=Pt(40),
                bold=True, color=BLUE_ACCENT, align=PP_ALIGN.CENTER)

    # Body lines — FONT_BODY (sans)
    lines = [
        ("🎱  Built a working 2D billiards simulation using Python and Pygame", LIGHT),
        ("📐  Used the Euler method to update ball positions each frame",        LIGHT),
        ("💥  Applied Newton's Law of Restitution for ball-ball collisions",     LIGHT),
        ("🌊  Added damping to simulate friction and energy loss over time",     LIGHT),
        ("📊  Ran 3 experiments — results confirmed the physics principles",     GREEN),
        ("✅  All project goals from the proposal were successfully completed",   GREEN),
    ]
    y = Inches(2.15)
    for text, col in lines:
        add_textbox(slide, text,
                    Inches(1.5), y, Inches(10.3), Inches(0.48),
                    font_name=FONT_BODY, font_size=Pt(18), color=col)
        y += Inches(0.5)

    add_textbox(slide, "Thank You  —  Questions?",
                Inches(1.3), Inches(6.05), Inches(10.7), Inches(0.5),
                font_name=FONT_TITLE, font_size=Pt(22),
                bold=True, color=DIM, align=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────────────────
    out = "Billiards_Simulation_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")
    print("Fonts used:")
    print(f"  Titles  : {FONT_TITLE}  (serif — Georgia)")
    print(f"  Body    : {FONT_BODY}  (sans-serif — Calibri Light)")
    print(f"  Formulas: {FONT_CODE}  (monospace — Courier New)")


if __name__ == "__main__":
    build()
